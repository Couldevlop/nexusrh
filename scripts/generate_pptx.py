"""
NexusRH — Générateur de présentations PPTX
- Document 1 : Présentation commerciale (apporteur d'affaires / commercial)
- Document 2 : Spécifications fonctionnelles techniques
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import os

# ─── PALETTE DE COULEURS PREMIUM ───────────────────────────────────────────────
INDIGO_900  = RGBColor(0x1E, 0x1B, 0x4B)   # fond sombre
INDIGO_800  = RGBColor(0x31, 0x2E, 0x81)   # fond slides
INDIGO_600  = RGBColor(0x4F, 0x46, 0xE5)   # primaire NexusRH
INDIGO_400  = RGBColor(0x81, 0x8C, 0xF8)   # accent clair
INDIGO_100  = RGBColor(0xE0, 0xE7, 0xFF)   # très clair
VIOLET_500  = RGBColor(0x8B, 0x5C, 0xF6)   # accent violet
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_100    = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_300    = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_500    = RGBColor(0x6B, 0x72, 0x80)
GRAY_700    = RGBColor(0x37, 0x41, 0x51)
GREEN_400   = RGBColor(0x34, 0xD3, 0x99)
AMBER_400   = RGBColor(0xFB, 0xBF, 0x24)
RED_400     = RGBColor(0xF8, 0x71, 0x71)
CYAN_400    = RGBColor(0x22, 0xD3, 0xEE)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── HELPERS ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # totalement vide
    return prs.slides.add_slide(layout)

def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txb

def add_paragraph(tf, text, font_size=14, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, italic=False, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return p

def add_bullet_box(slide, items, left, top, width, height,
                   font_size=14, color=WHITE, bullet="▸  ", title=None, title_color=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(font_size + 2)
        run.font.bold = True
        run.font.color.rgb = title_color or INDIGO_400
        run.font.name = "Segoe UI"
        first = False
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}{item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return txb

def card(slide, left, top, width, height, bg_color, title, title_color,
         body_lines, body_color=WHITE, title_size=16, body_size=13, radius_hint=None):
    """Card avec fond coloré, titre et lignes de body."""
    add_rect(slide, left, top, width, height, bg_color)
    # accent bar top
    add_rect(slide, left, top, width, Inches(0.06), title_color)
    # title
    add_text(slide, title, left + Inches(0.2), top + Inches(0.12),
             width - Inches(0.4), Inches(0.4),
             font_size=title_size, bold=True, color=title_color)
    # body
    y = top + Inches(0.55)
    for line in body_lines:
        add_text(slide, line, left + Inches(0.2), y,
                 width - Inches(0.4), Inches(0.35),
                 font_size=body_size, color=body_color, wrap=True)
        y += Inches(0.32)

def divider(slide, top, color=INDIGO_600):
    add_rect(slide, Inches(0.5), top, SLIDE_W - Inches(1), Pt(1.5), color)

def slide_header(slide, title, subtitle=None,
                 title_color=WHITE, sub_color=INDIGO_400, accent=True):
    if accent:
        add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), INDIGO_600)
    add_text(slide, title, Inches(0.5), Inches(0.18), SLIDE_W - Inches(1), Inches(0.7),
             font_size=28, bold=True, color=title_color)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(0.88), SLIDE_W - Inches(1), Inches(0.4),
                 font_size=15, color=sub_color, italic=True)

def footer(slide, text="NexusRH — SIRH SaaS Multi-Tenant  |  Confidentiel"):
    add_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), INDIGO_900)
    add_text(slide, text, Inches(0.3), SLIDE_H - Inches(0.33),
             SLIDE_W - Inches(0.6), Inches(0.3),
             font_size=9, color=GRAY_300, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  DOCUMENT 1 — PRÉSENTATION COMMERCIALE
# ═══════════════════════════════════════════════════════════════════════════════

def build_commercial_deck(output_path: str):
    prs = new_prs()

    # ── SLIDE 1 : COVER ──────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    # gradient bar gauche
    add_rect(sl, 0, 0, Inches(0.5), SLIDE_H, INDIGO_600)
    # accent lines décoratives
    add_rect(sl, Inches(0.5), Inches(2.2), SLIDE_W - Inches(0.5), Pt(2), INDIGO_600)
    add_rect(sl, Inches(0.5), Inches(5.1), SLIDE_W - Inches(0.5), Pt(1), VIOLET_500)
    # logo placeholder
    add_rect(sl, Inches(0.9), Inches(0.6), Inches(1.2), Inches(1.2), INDIGO_600)
    add_text(sl, "NRH", Inches(0.9), Inches(0.65), Inches(1.2), Inches(1.1),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # titre
    add_text(sl, "NexusRH", Inches(2.4), Inches(0.55), Inches(9), Inches(1.1),
             font_size=52, bold=True, color=WHITE)
    add_text(sl, "Le SIRH SaaS multi-tenant de nouvelle génération",
             Inches(2.4), Inches(1.65), Inches(9), Inches(0.7),
             font_size=20, color=INDIGO_400, italic=True)
    # séparateur
    add_rect(sl, Inches(0.9), Inches(2.5), Inches(10.5), Pt(2), INDIGO_600)
    # tagline
    add_text(sl, "Pilotez vos RH. Libérez vos équipes. Décidez avec l'IA.",
             Inches(0.9), Inches(2.75), Inches(11), Inches(0.8),
             font_size=22, bold=True, color=WHITE)
    # pills
    for i, (label, col) in enumerate([
        ("Multi-tenant SaaS", INDIGO_600),
        ("IA embarquée", VIOLET_500),
        ("Conformité France 2024", GREEN_400),
        ("Open API", CYAN_400),
    ]):
        x = Inches(0.9) + i * Inches(2.9)
        add_rect(sl, x, Inches(3.5), Inches(2.7), Inches(0.42), col)
        add_text(sl, label, x + Inches(0.1), Inches(3.5),
                 Inches(2.5), Inches(0.42),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # bas
    add_text(sl, "Présentation commerciale  ·  Confidentiel  ·  2025",
             Inches(0.9), Inches(5.2), Inches(10), Inches(0.4),
             font_size=12, color=GRAY_300, italic=True)
    add_text(sl, "Propulsé par Claude AI (Anthropic) & Mistral AI",
             Inches(0.9), Inches(5.6), Inches(10), Inches(0.4),
             font_size=11, color=GRAY_500, italic=True)

    # ── SLIDE 2 : PROBLÈME ───────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    slide_header(sl, "Le défi RH des entreprises aujourd'hui",
                 "Pourquoi les outils actuels ne suffisent plus")
    footer(sl)

    problems = [
        ("🗂  Silos de données", "Paie, absences, recrutement dans des outils disparates → erreurs, double saisie, perte de temps."),
        ("🔒  Sécurité & conformité", "RGPD, DSN, bulletins légaux : la complexité réglementaire française dépasse les outils génériques."),
        ("💸  Coûts cachés", "Licences multiples, intégrations coûteuses, consultants externes pour chaque module."),
        ("🤖  Zéro intelligence", "Aucune anticipation : turnover, burn-out, pics d'absentéisme détectés trop tard."),
        ("🏢  Croissance bloquée", "Les ESN et cabinets RH ne peuvent pas gérer plusieurs clients sur une même plateforme sécurisée."),
        ("📊  Reporting limité", "Tableaux de bord statiques, exports Excel manuels, aucune vision temps réel."),
    ]
    cols = [Inches(0.5), Inches(7.0)]
    for idx, (title, desc) in enumerate(problems):
        col = cols[idx % 2]
        row = idx // 2
        y = Inches(1.5) + row * Inches(1.7)
        bg = INDIGO_800 if idx % 2 == 0 else RGBColor(0x2D, 0x2A, 0x7A)
        add_rect(sl, col, y, Inches(6.1), Inches(1.55), bg)
        add_rect(sl, col, y, Inches(0.06), Inches(1.55), RED_400)
        add_text(sl, title, col + Inches(0.15), y + Inches(0.1),
                 Inches(5.8), Inches(0.45), font_size=15, bold=True, color=WHITE)
        add_text(sl, desc, col + Inches(0.15), y + Inches(0.55),
                 Inches(5.8), Inches(0.85), font_size=11, color=GRAY_300, wrap=True)

    # ── SLIDE 3 : SOLUTION ───────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    slide_header(sl, "NexusRH — Une plateforme. Tous vos besoins RH.",
                 "Architecture SaaS multi-tenant, conformité France, IA de dernière génération")
    footer(sl)

    add_rect(sl, Inches(0.5), Inches(1.5), SLIDE_W - Inches(1), Inches(5.5), INDIGO_800)
    add_rect(sl, Inches(0.5), Inches(1.5), Inches(0.08), Inches(5.5), INDIGO_600)

    points = [
        ("✅  Une seule plateforme", "Tous les modules RH intégrés nativement — paie, absences, recrutement, formation, carrière, frais."),
        ("✅  Multi-tenant sécurisé", "Chaque client dispose de son propre schéma PostgreSQL isolé. Les données ne se mélangent jamais."),
        ("✅  IA embarquée", "Scoring rétention, génération de documents RH, analyse CV, comparateur Claude vs Mistral pour le sourcing."),
        ("✅  Conformité France 2024", "Moteur de paie SYNTEC/BTP, DSN, RGPD (NIR + IBAN chiffrés AES-256), bulletins PDF légaux."),
        ("✅  White-label & SaaS", "Thématisation dynamique par client (logo, couleurs). Interface employee self-service intégrée."),
        ("✅  API ouverte", "API REST documentée (Swagger), webhooks, intégrations tierces (comptabilité, ATS externe, SI paie)."),
    ]
    for i, (title, desc) in enumerate(points):
        y = Inches(1.65) + i * Inches(0.83)
        add_text(sl, title, Inches(0.85), y, Inches(3.5), Inches(0.5),
                 font_size=13, bold=True, color=GREEN_400)
        add_text(sl, desc, Inches(4.5), y, Inches(8.5), Inches(0.7),
                 font_size=12, color=GRAY_300, wrap=True)

    # ── SLIDE 4 : MODULES ────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    slide_header(sl, "Modules fonctionnels intégrés",
                 "De l'embauche à la retraite — tout en un")
    footer(sl)

    modules = [
        ("👥  Employés", INDIGO_600, ["Dossier complet", "Import CSV", "Documents", "RGPD"]),
        ("💰  Paie", VIOLET_500, ["Moteur calcul", "Bulletins PDF", "DSN", "SEPA"]),
        ("🗓  Absences", GREEN_400, ["Workflow multi-niveaux", "Soldes temps réel", "Planning"]),
        ("🎯  Recrutement", AMBER_400, ["Pipeline kanban", "Analyse CV IA", "Sourcing IA"]),
        ("📚  Formation", CYAN_400, ["Catalogue", "Inscriptions", "Attestations PDF"]),
        ("🧾  Notes de frais", RED_400, ["Lignes détaillées", "Justificatifs S3", "Workflow"]),
        ("📈  Carrière", RGBColor(0xA7, 0x8B, 0xFA), ["Compétences", "9-box", "Entretiens"]),
        ("📊  Reporting IA", RGBColor(0x5E, 0xEA, 0xD4), ["KPIs temps réel", "Prédictif", "Exports"]),
    ]
    positions = [
        (Inches(0.4),  Inches(1.5)),
        (Inches(3.35), Inches(1.5)),
        (Inches(6.3),  Inches(1.5)),
        (Inches(9.25), Inches(1.5)),
        (Inches(0.4),  Inches(4.0)),
        (Inches(3.35), Inches(4.0)),
        (Inches(6.3),  Inches(4.0)),
        (Inches(9.25), Inches(4.0)),
    ]
    for (title, accent_col, items), (x, y) in zip(modules, positions):
        add_rect(sl, x, y, Inches(2.7), Inches(2.2), INDIGO_800)
        add_rect(sl, x, y, Inches(2.7), Inches(0.07), accent_col)
        add_text(sl, title, x + Inches(0.12), y + Inches(0.12),
                 Inches(2.46), Inches(0.5), font_size=13, bold=True, color=WHITE)
        for j, item in enumerate(items):
            add_text(sl, f"  • {item}", x + Inches(0.12), y + Inches(0.65) + j * Inches(0.35),
                     Inches(2.46), Inches(0.35), font_size=11, color=GRAY_300)

    # ── SLIDE 5 : MULTI-TENANT ───────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    slide_header(sl, "Architecture multi-tenant — Isolation totale des données",
                 "Chaque client = son propre schéma PostgreSQL isolé")
    footer(sl)

    # Schéma visuel simplifié
    add_rect(sl, Inches(0.5), Inches(1.5), SLIDE_W - Inches(1), Inches(1.1),
             RGBColor(0x1A, 0x17, 0x4E))
    add_text(sl, "🌐  Plateforme NexusRH (schema : platform)",
             Inches(0.7), Inches(1.6), Inches(11), Inches(0.5),
             font_size=14, bold=True, color=INDIGO_400, align=PP_ALIGN.CENTER)
    add_text(sl, "Gestion des tenants  ·  Authentification  ·  Feature flags  ·  Logs globaux",
             Inches(0.7), Inches(1.95), Inches(11), Inches(0.4),
             font_size=11, color=GRAY_300, align=PP_ALIGN.CENTER)

    tenants = [
        ("TechCorp SAS", INDIGO_600, "#4F46E5", "Pro · 50 emp."),
        ("Artisan Pro SARL", GREEN_400, "#16A34A", "Starter · 18 emp."),
        ("CabinetRH & Co", VIOLET_500, "#8B5CF6", "Enterprise · 200 emp."),
        ("StartupX", AMBER_400, "#F59E0B", "Trial · 5 emp."),
    ]
    for i, (name, col, hex_col, info) in enumerate(tenants):
        x = Inches(0.5) + i * Inches(3.1)
        add_rect(sl, x, Inches(3.0), Inches(2.9), Inches(3.0), INDIGO_800)
        add_rect(sl, x, Inches(3.0), Inches(2.9), Inches(0.07), col)
        add_rect(sl, x + Inches(0.1), Inches(3.15), Inches(0.5), Inches(0.5), col)
        add_text(sl, name[:2].upper(), x + Inches(0.1), Inches(3.15),
                 Inches(0.5), Inches(0.5), font_size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(sl, name, x + Inches(0.7), Inches(3.17),
                 Inches(2.1), Inches(0.45), font_size=13, bold=True, color=WHITE)
        add_text(sl, info, x + Inches(0.7), Inches(3.6),
                 Inches(2.1), Inches(0.35), font_size=10, color=col)
        items = ["Employés isolés", "Paie propre", "Config custom", f"Thème {hex_col}"]
        for j, it in enumerate(items):
            add_text(sl, f"  • {it}", x + Inches(0.15), Inches(4.15) + j * Inches(0.4),
                     Inches(2.6), Inches(0.38), font_size=10.5, color=GRAY_300)

    add_text(sl, "🔒  Isolation totale PostgreSQL  —  SET search_path par requête  —  Zéro fuite de données inter-tenant",
             Inches(0.5), Inches(6.2), SLIDE_W - Inches(1), Inches(0.4),
             font_size=11, bold=True, color=GREEN_400, align=PP_ALIGN.CENTER)

    # ── SLIDE 6 : CIBLES ─────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    slide_header(sl, "Pour qui est fait NexusRH ?",
                 "Trois profils cibles, une solution unique")
    footer(sl)

    targets = [
        (
            "🏢  PME & ETI (50–2 000 salariés)",
            INDIGO_600,
            [
                "Besoin : remplacer plusieurs outils silos (Silae, Sage, Cegid)",
                "Valeur : one-stop-shop RH, conformité France clé en main",
                "ROI : -40% coûts logiciels, -60% saisies manuelles",
                "Décideur : DRH, DAF, DSI",
            ]
        ),
        (
            "🤝  Cabinets RH & ESN",
            VIOLET_500,
            [
                "Besoin : gérer plusieurs clients sur une plateforme unique",
                "Valeur : multi-tenant natif, white-label, marque propre",
                "ROI : un seul abonnement pour N clients, facturation mutualisée",
                "Décideur : Associé, Directeur pôle social",
            ]
        ),
        (
            "🌍  Groupes multi-entités",
            GREEN_400,
            [
                "Besoin : consolidation RH cross-filiales, reporting groupe",
                "Valeur : isolation légale + vision consolidée, SEPA multi-SIRET",
                "ROI : gouvernance RH unifiée, conformité multi-CCN",
                "Décideur : DRH Groupe, DOSI",
            ]
        ),
    ]
    for i, (title, col, items) in enumerate(targets):
        x = Inches(0.4) + i * Inches(4.25)
        add_rect(sl, x, Inches(1.5), Inches(4.05), Inches(5.0), INDIGO_800)
        add_rect(sl, x, Inches(1.5), Inches(4.05), Inches(0.07), col)
        add_text(sl, title, x + Inches(0.15), Inches(1.6),
                 Inches(3.75), Inches(0.65), font_size=14, bold=True, color=WHITE, wrap=True)
        for j, item in enumerate(items):
            add_text(sl, f"  • {item}", x + Inches(0.15), Inches(2.35) + j * Inches(0.8),
                     Inches(3.75), Inches(0.75), font_size=12, color=GRAY_300, wrap=True)

    # ── SLIDE 7 : IA ─────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    slide_header(sl, "Intelligence Artificielle — Un avantage concurrentiel majeur",
                 "Claude (Anthropic) + Mistral pour le meilleur des deux mondes")
    footer(sl)

    # Claude block
    add_rect(sl, Inches(0.5), Inches(1.55), Inches(5.9), Inches(4.9), INDIGO_800)
    add_rect(sl, Inches(0.5), Inches(1.55), Inches(0.07), Inches(4.9), INDIGO_600)
    add_text(sl, "🤖  Claude (Anthropic)", Inches(0.7), Inches(1.65),
             Inches(5.5), Inches(0.55), font_size=16, bold=True, color=INDIGO_400)
    claude_items = [
        "Assistant RH conversationnel (streaming SSE)",
        "Génération de documents légaux complets",
        "   (CDI, CDD, lettre avertissement, rupture conventionnelle…)",
        "Analyse CV candidates (scoring + synthèse)",
        "Scoring rétention & risque burn-out (nuit)",
        "Réponses en français + références art. Code du travail",
        "Modèle : claude-sonnet-4-20250514",
    ]
    for j, item in enumerate(claude_items):
        col = GRAY_300 if not item.startswith("   ") else GRAY_500
        add_text(sl, f"  • {item.strip()}" if not item.startswith("   ") else item,
                 Inches(0.7), Inches(2.3) + j * Inches(0.46),
                 Inches(5.5), Inches(0.44), font_size=11.5, color=col, wrap=True)

    # Mistral block
    add_rect(sl, Inches(6.9), Inches(1.55), Inches(5.9), Inches(4.9), INDIGO_800)
    add_rect(sl, Inches(6.9), Inches(1.55), Inches(0.07), Inches(4.9), VIOLET_500)
    add_text(sl, "🧠  Mistral AI", Inches(7.1), Inches(1.65),
             Inches(5.5), Inches(0.55), font_size=16, bold=True, color=VIOLET_500)
    mistral_items = [
        "Sourcing candidats (stratégies multi-plateformes)",
        "Modèle souverain européen (conformité données)",
        "Analyse comparative automatique Claude vs Mistral",
        "   → résultat le plus riche sélectionné automatiquement",
        "Latence & coût optimisés pour les traitements batch",
        "Modèle : mistral-large-latest",
    ]
    for j, item in enumerate(mistral_items):
        col = GRAY_300 if not item.startswith("   ") else GRAY_500
        add_text(sl, f"  • {item.strip()}" if not item.startswith("   ") else item,
                 Inches(7.1), Inches(2.3) + j * Inches(0.46),
                 Inches(5.5), Inches(0.44), font_size=11.5, color=col, wrap=True)

    # badge comparateur
    add_rect(sl, Inches(4.6), Inches(3.6), Inches(4.05), Inches(0.9), INDIGO_600)
    add_text(sl, "⚡  Comparateur automatique — le meilleur résultat gagne",
             Inches(4.6), Inches(3.62), Inches(4.05), Inches(0.85),
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

    # ── SLIDE 8 : PLANS & TARIFICATION ───────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    slide_header(sl, "Plans & tarification SaaS",
                 "Facturation mensuelle par utilisateur actif — sans engagement minimal")
    footer(sl)

    plans = [
        ("🌱  Trial", GRAY_500, "Gratuit 30 jours", [
            "Jusqu'à 10 utilisateurs",
            "20 employés max",
            "Tous les modules",
            "Support email",
        ], False),
        ("🚀  Starter", GREEN_400, "Dès 299 €/mois", [
            "50 utilisateurs",
            "100 employés",
            "Modules core RH",
            "IA basique",
            "Support standard",
        ], False),
        ("💼  Pro", INDIGO_600, "Dès 799 €/mois", [
            "200 utilisateurs",
            "500 employés",
            "Tous les modules",
            "IA complète (Claude + Mistral)",
            "API + webhooks",
            "Support prioritaire",
        ], True),
        ("🏛  Enterprise", AMBER_400, "Sur devis", [
            "Illimité",
            "Multi-entités",
            "White-label",
            "SLA 99,9%",
            "Déploiement on-premise",
            "CSM dédié",
        ], False),
    ]
    for i, (title, col, price, items, highlighted) in enumerate(plans):
        x = Inches(0.4) + i * Inches(3.2)
        bg = INDIGO_600 if highlighted else INDIGO_800
        add_rect(sl, x, Inches(1.55), Inches(3.0), Inches(5.0), bg)
        add_rect(sl, x, Inches(1.55), Inches(3.0), Inches(0.07), col)
        if highlighted:
            add_rect(sl, x - Inches(0.05), Inches(1.4), Inches(3.1), Inches(0.22),
                     AMBER_400)
            add_text(sl, "⭐  PLUS POPULAIRE", x, Inches(1.4),
                     Inches(3.0), Inches(0.22), font_size=10, bold=True,
                     color=INDIGO_900, align=PP_ALIGN.CENTER)
        add_text(sl, title, x + Inches(0.15), Inches(1.7),
                 Inches(2.7), Inches(0.55), font_size=15, bold=True, color=WHITE)
        add_text(sl, price, x + Inches(0.15), Inches(2.25),
                 Inches(2.7), Inches(0.5), font_size=18, bold=True, color=col)
        for j, item in enumerate(items):
            add_text(sl, f"  ✓  {item}", x + Inches(0.15), Inches(2.85) + j * Inches(0.52),
                     Inches(2.7), Inches(0.5), font_size=11, color=GRAY_300, wrap=True)

    add_text(sl, "🔧  Module IA, signature électronique, SSO SAML en option sur Starter · Offre cabinet RH disponible (multi-tenant partenaire)",
             Inches(0.4), Inches(6.8), SLIDE_W - Inches(0.8), Inches(0.35),
             font_size=10, color=GRAY_500, italic=True, align=PP_ALIGN.CENTER)

    # ── SLIDE 9 : ROI ────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    slide_header(sl, "ROI & bénéfices mesurables",
                 "Des gains concrets, quantifiables dès le premier trimestre")
    footer(sl)

    kpis = [
        ("⏱  -60%", "de temps de saisie\npour les équipes RH", GREEN_400),
        ("💰  -40%", "de coût logiciel\nvs solution multi-outils", INDIGO_400),
        ("🎯  +35%", "de satisfaction\ndes managers terrain", AMBER_400),
        ("🔍  -25%", "de turnover\ngrâce au scoring IA", VIOLET_500),
    ]
    for i, (num, label, col) in enumerate(kpis):
        x = Inches(0.5) + i * Inches(3.15)
        add_rect(sl, x, Inches(1.5), Inches(2.95), Inches(2.2), INDIGO_800)
        add_rect(sl, x, Inches(1.5), Inches(2.95), Inches(0.07), col)
        add_text(sl, num, x + Inches(0.15), Inches(1.65),
                 Inches(2.65), Inches(0.9), font_size=34, bold=True, color=col)
        add_text(sl, label, x + Inches(0.15), Inches(2.55),
                 Inches(2.65), Inches(0.9), font_size=12, color=GRAY_300, wrap=True)

    # Use cases
    add_text(sl, "📋  Cas clients types", Inches(0.5), Inches(3.9),
             Inches(12), Inches(0.5), font_size=16, bold=True, color=WHITE)
    cases = [
        ("PME 80 salariés — Industrie", "Remplace Silae + Cegid + fichiers Excel. Économie : 18 000 €/an. Paie clôturée en 2h au lieu de 2 jours."),
        ("Cabinet RH — 12 clients", "1 seule plateforme NexusRH gère 12 dossiers. Temps de paramétrage nouvel client : 15 min. Marge doublée."),
        ("ETI 400 salariés — Services", "Turnover réduit de 22% en 6 mois grâce aux alertes IA. Absentéisme détecté 3 semaines avant le pic."),
    ]
    for i, (title, desc) in enumerate(cases):
        y = Inches(4.5) + i * Inches(0.75)
        add_rect(sl, Inches(0.5), y, Inches(12.3), Inches(0.68), INDIGO_800)
        add_rect(sl, Inches(0.5), y, Inches(0.06), Inches(0.68), GREEN_400)
        add_text(sl, title, Inches(0.7), y + Inches(0.05),
                 Inches(3.5), Inches(0.55), font_size=12, bold=True, color=WHITE)
        add_text(sl, desc, Inches(4.2), y + Inches(0.05),
                 Inches(8.5), Inches(0.55), font_size=11, color=GRAY_300, wrap=True)

    # ── SLIDE 10 : ROADMAP ───────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    slide_header(sl, "Roadmap produit 2025–2026",
                 "Vision ambitieuse, livraisons régulières")
    footer(sl)

    phases = [
        ("Q2 2025\n✅ Disponible", INDIGO_600, [
            "Core RH complet (50+ endpoints)",
            "Moteur de paie SYNTEC",
            "Espace employé self-service",
            "IA : Claude + Mistral",
            "Multi-tenant production",
        ]),
        ("Q3 2025\n🔄 En cours", AMBER_400, [
            "Signature électronique (DocuSign)",
            "DSN automatique URSSAF",
            "Mobile App (iOS/Android)",
            "Portail client white-label",
            "SSO SAML 2.0 / OIDC",
        ]),
        ("Q4 2025\n📅 Planifié", VIOLET_500, [
            "Analytics prédictifs avancés",
            "Multi-pays (BE, CH, LU)",
            "Marketplace intégrations",
            "Chatbot HR Copilot",
            "Kiosk mode (pointeuse)",
        ]),
        ("2026\n🌟 Vision", GREEN_400, [
            "IA générative bulletins",
            "Benchmark salarial IA",
            "HRIS API for developers",
            "Acquisition & onboarding IA",
            "Conformité DORA/NIS2",
        ]),
    ]
    for i, (period, col, items) in enumerate(phases):
        x = Inches(0.4) + i * Inches(3.15)
        add_rect(sl, x, Inches(1.55), Inches(3.0), Inches(5.0), INDIGO_800)
        add_rect(sl, x, Inches(1.55), Inches(3.0), Inches(0.07), col)
        add_text(sl, period, x + Inches(0.15), Inches(1.65),
                 Inches(2.7), Inches(0.75), font_size=13, bold=True, color=col, wrap=True)
        for j, item in enumerate(items):
            add_text(sl, f"  • {item}", x + Inches(0.15), Inches(2.55) + j * Inches(0.48),
                     Inches(2.7), Inches(0.46), font_size=11, color=GRAY_300, wrap=True)

    # ── SLIDE 11 : POURQUOI NOUS ─────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    slide_header(sl, "Pourquoi choisir NexusRH ?",
                 "5 raisons décisives vs la concurrence")
    footer(sl)

    reasons = [
        ("1️⃣  Nativement multi-tenant", INDIGO_600,
         "Pas un ajout tardif : conçu dès l'origine pour gérer N clients en isolation totale. "
         "Vos données ne cohabitent JAMAIS avec celles d'un autre client."),
        ("2️⃣  IA de pointe, pas du marketing", VIOLET_500,
         "Deux modèles IA (Claude + Mistral) intégrés en production avec comparateur automatique. "
         "Le meilleur résultat est sélectionné objectivement par score de richesse."),
        ("3️⃣  Conformité France clé en main", GREEN_400,
         "Moteur de paie SYNTEC/BTP, bulletins PDF légaux, DSN, RGPD (NIR + IBAN AES-256), "
         "50 rubriques de cotisations France 2024 préconfigurées."),
        ("4️⃣  Open-source friendly, API first", CYAN_400,
         "API REST documentée Swagger, webhooks, SDK TypeScript open-source. "
         "Votre DSI peut connecter n'importe quel outil existant en quelques heures."),
        ("5️⃣  Déploiement en 15 minutes", AMBER_400,
         "Création d'un nouveau tenant client : formulaire 3 étapes → schéma provisionné → "
         "admin notifié par email. Zéro intervention technique."),
    ]
    for i, (title, col, desc) in enumerate(reasons):
        y = Inches(1.5) + i * Inches(1.0)
        add_rect(sl, Inches(0.5), y, SLIDE_W - Inches(1), Inches(0.9), INDIGO_800)
        add_rect(sl, Inches(0.5), y, Inches(0.07), Inches(0.9), col)
        add_text(sl, title, Inches(0.72), y + Inches(0.07),
                 Inches(3.8), Inches(0.75), font_size=13, bold=True, color=col)
        add_text(sl, desc, Inches(4.7), y + Inches(0.07),
                 Inches(8.4), Inches(0.75), font_size=11.5, color=GRAY_300, wrap=True)

    # ── SLIDE 12 : CTA / CONTACT ─────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO_900)
    add_rect(sl, 0, 0, Inches(0.5), SLIDE_H, INDIGO_600)
    add_rect(sl, Inches(0.5), SLIDE_H / 2 - Pt(1), SLIDE_W - Inches(0.5), Pt(2), INDIGO_600)

    add_text(sl, "Prêt à transformer vos RH ?",
             Inches(1.2), Inches(1.0), Inches(11), Inches(1.1),
             font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "Demandez une démo live — disponible sous 24h",
             Inches(1.2), Inches(2.1), Inches(11), Inches(0.7),
             font_size=20, color=INDIGO_400, italic=True, align=PP_ALIGN.CENTER)

    for i, (icon, label, val, col) in enumerate([
        ("📧", "Email commercial", "contact@nexusrh.com", INDIGO_400),
        ("🌐", "Site web", "www.nexusrh.com", GREEN_400),
        ("📞", "Téléphone", "+33 1 XX XX XX XX", AMBER_400),
        ("💬", "Demo live", "Calendly — 30 min", VIOLET_500),
    ]):
        x = Inches(1.0) + i * Inches(3.0)
        add_rect(sl, x, Inches(3.2), Inches(2.8), Inches(1.6), INDIGO_800)
        add_rect(sl, x, Inches(3.2), Inches(2.8), Inches(0.06), col)
        add_text(sl, f"{icon}  {label}", x + Inches(0.15), Inches(3.3),
                 Inches(2.5), Inches(0.5), font_size=12, bold=True, color=col)
        add_text(sl, val, x + Inches(0.15), Inches(3.78),
                 Inches(2.5), Inches(0.5), font_size=11, color=WHITE)

    add_text(sl, "NexusRH — SIRH SaaS Multi-Tenant  ·  Propulsé par Claude AI & Mistral AI",
             Inches(1.2), Inches(5.2), Inches(11), Inches(0.4),
             font_size=11, color=GRAY_500, italic=True, align=PP_ALIGN.CENTER)
    add_text(sl, "© 2025 OpenLab — Document confidentiel. Ne pas diffuser sans autorisation.",
             Inches(1.2), Inches(5.6), Inches(11), Inches(0.4),
             font_size=10, color=GRAY_700, italic=True, align=PP_ALIGN.CENTER)
    footer(sl)

    prs.save(output_path)
    print(f"[OK] Document commercial sauvegarde : {output_path}")


# ═══════════════════════════════════════════════════════════════════════════════
#  DOCUMENT 2 — SPÉCIFICATIONS FONCTIONNELLES TECHNIQUES
# ═══════════════════════════════════════════════════════════════════════════════

TECH_BG    = RGBColor(0x0F, 0x17, 0x2A)   # très sombre bleu-marine
TECH_PANEL = RGBColor(0x1A, 0x27, 0x40)   # panels
TECH_ACCENT= RGBColor(0x06, 0xB6, 0xD4)   # cyan tech
TECH_GREEN = RGBColor(0x10, 0xB9, 0x81)
TECH_YELLOW= RGBColor(0xF5, 0x9E, 0x0B)
TECH_RED   = RGBColor(0xEF, 0x44, 0x44)
TECH_PURPLE= RGBColor(0x7C, 0x3A, 0xED)
TECH_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
TECH_GRAY  = RGBColor(0x94, 0xA3, 0xB8)
CODE_BG    = RGBColor(0x0D, 0x1B, 0x2A)

def tech_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), TECH_ACCENT)
    add_text(slide, title, Inches(0.5), Inches(0.18), SLIDE_W - Inches(1), Inches(0.7),
             font_size=26, bold=True, color=TECH_WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(0.82), SLIDE_W - Inches(1), Inches(0.38),
                 font_size=13, color=TECH_ACCENT, italic=True)

def tech_footer(slide, page_info=""):
    add_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), RGBColor(0x08, 0x0F, 0x1D))
    add_text(slide, "NexusRH — Spécifications Fonctionnelles Techniques  |  Confidentiel — Diffusion restreinte",
             Inches(0.3), SLIDE_H - Inches(0.34), Inches(10), Inches(0.3),
             font_size=9, color=TECH_GRAY, italic=True)
    if page_info:
        add_text(slide, page_info, Inches(11.5), SLIDE_H - Inches(0.34),
                 Inches(1.5), Inches(0.3), font_size=9, color=TECH_GRAY,
                 align=PP_ALIGN.RIGHT)

def code_block(slide, code_lines, left, top, width, height, title=None):
    add_rect(slide, left, top, width, height, CODE_BG)
    add_rect(slide, left, top, width, Pt(1.5), TECH_ACCENT)
    if title:
        add_text(slide, title, left + Inches(0.12), top + Inches(0.08),
                 width - Inches(0.24), Inches(0.35),
                 font_size=10, bold=True, color=TECH_ACCENT)
        offset = Inches(0.42)
    else:
        offset = Inches(0.1)
    txb = slide.shapes.add_textbox(left + Inches(0.12), top + offset,
                                   width - Inches(0.24), height - offset - Inches(0.1))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in code_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9.5)
        run.font.name = "Consolas"
        run.font.color.rgb = TECH_WHITE

def build_technical_spec(output_path: str):
    prs = new_prs()

    # ── SLIDE 1 : COVER TECH ─────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    # lignes décoratives
    for i in range(8):
        y = Inches(0.6) + i * Inches(0.85)
        opacity_col = RGBColor(
            int(0x06 * (0.15 + 0.05 * i)),
            int(0xB6 * (0.15 + 0.05 * i)),
            int(0xD4 * (0.15 + 0.05 * i)),
        )
        add_rect(sl, 0, y, SLIDE_W, Pt(0.5), RGBColor(0x0F, 0x2D, 0x40))
    add_rect(sl, 0, 0, Inches(0.06), SLIDE_H, TECH_ACCENT)
    add_rect(sl, Inches(0.5), Inches(1.9), SLIDE_W - Inches(0.5), Pt(2), TECH_ACCENT)
    add_rect(sl, Inches(0.5), Inches(4.9), SLIDE_W - Inches(0.5), Pt(1), TECH_PURPLE)

    add_rect(sl, Inches(0.8), Inches(0.6), Inches(1.2), Inches(1.2), TECH_ACCENT)
    add_text(sl, "NRH", Inches(0.8), Inches(0.65), Inches(1.2), Inches(1.1),
             font_size=28, bold=True, color=TECH_BG, align=PP_ALIGN.CENTER)

    add_text(sl, "NexusRH", Inches(2.3), Inches(0.55), Inches(10), Inches(1.1),
             font_size=50, bold=True, color=TECH_WHITE)
    add_text(sl, "Spécifications Fonctionnelles Techniques",
             Inches(2.3), Inches(1.6), Inches(10), Inches(0.55),
             font_size=18, bold=True, color=TECH_ACCENT)
    add_text(sl, "SIRH SaaS Multi-Tenant — Architecture, API, Sécurité, IA",
             Inches(0.8), Inches(2.15), Inches(11.5), Inches(0.5),
             font_size=14, color=TECH_GRAY, italic=True)

    meta = [
        ("Version", "1.0.0 — Avril 2025"),
        ("Stack", "Node.js 20 · Fastify 4 · Drizzle ORM · React 18 · PostgreSQL 16"),
        ("IA", "Claude Sonnet 4 (Anthropic) · Mistral Large (Mistral AI)"),
        ("Auteur", "OpenLab — Équipe Architecture"),
        ("Diffusion", "Interne + Partenaires techniques autorisés"),
    ]
    for i, (k, v) in enumerate(meta):
        y = Inches(2.9) + i * Inches(0.42)
        add_text(sl, f"{k} :", Inches(0.8), y, Inches(1.8), Inches(0.4),
                 font_size=11, bold=True, color=TECH_ACCENT)
        add_text(sl, v, Inches(2.7), y, Inches(10), Inches(0.4),
                 font_size=11, color=TECH_WHITE)

    add_rect(sl, Inches(0.8), Inches(5.1), SLIDE_W - Inches(1.6), Inches(0.7), TECH_PANEL)
    add_text(sl,
             "⚠  Confidentiel — Ce document contient des informations techniques propriétaires. "
             "Ne pas divulguer sans accord écrit.",
             Inches(1.0), Inches(5.15), SLIDE_W - Inches(2), Inches(0.55),
             font_size=10, color=TECH_YELLOW, italic=True, wrap=True)

    # ── SLIDE 2 : SOMMAIRE ───────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "Table des matières")
    tech_footer(sl, "2 / 14")

    toc = [
        ("01", "Vue d'ensemble & architecture globale"),
        ("02", "Stack technique détaillée"),
        ("03", "Architecture multi-tenant (schema-per-tenant)"),
        ("04", "Authentification, autorisation & RBAC"),
        ("05", "Modules fonctionnels — API REST"),
        ("06", "Moteur de paie — Logique & calculs"),
        ("07", "Intégration IA — Claude & Mistral"),
        ("08", "Gestion des fichiers & stockage"),
        ("09", "Recherche & indexation (Meilisearch)"),
        ("10", "Sécurité & conformité RGPD"),
        ("11", "Infrastructure & déploiement"),
        ("12", "Performance & scalabilité"),
        ("13", "Tests & qualité"),
        ("14", "Annexes — Schémas & endpoints"),
    ]
    col_split = 7
    for i, (num, title) in enumerate(toc):
        col = 0 if i < col_split else 1
        row = i if i < col_split else i - col_split
        x = Inches(0.5) + col * Inches(6.5)
        y = Inches(1.45) + row * Inches(0.72)
        add_rect(sl, x, y, Inches(6.2), Inches(0.62), TECH_PANEL)
        add_rect(sl, x, y, Inches(0.55), Inches(0.62), TECH_ACCENT)
        add_text(sl, num, x, y, Inches(0.55), Inches(0.62),
                 font_size=13, bold=True, color=TECH_BG, align=PP_ALIGN.CENTER)
        add_text(sl, title, x + Inches(0.65), y + Inches(0.1),
                 Inches(5.4), Inches(0.42), font_size=12, color=TECH_WHITE)

    # ── SLIDE 3 : VUE D'ENSEMBLE ─────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "01 — Vue d'ensemble & architecture globale",
                "Système distribué, modulaire, cloud-native")
    tech_footer(sl, "3 / 14")

    add_rect(sl, Inches(0.5), Inches(1.45), SLIDE_W - Inches(1), Inches(4.9), TECH_PANEL)
    # layers
    layers = [
        ("Frontend Layer", TECH_ACCENT, "React 18 + Vite 5 · shadcn/ui · TanStack Query 5 · Zustand 4 · React Router 6"),
        ("API Gateway", TECH_GREEN, "Fastify 4 · JWT + OAuth2 · Rate limiting · CORS · Swagger UI · Multipart"),
        ("Business Logic", TECH_PURPLE, "Modules RH · Moteur paie · Workflow absences/frais · RBAC middleware"),
        ("AI Layer", TECH_YELLOW, "Claude Sonnet 4 (SSE streaming) · Mistral Large · Comparateur automatique"),
        ("Data Layer", TECH_RED, "PostgreSQL 16 (multi-schema) · Drizzle ORM · Redis 7 (sessions/cache)"),
        ("Infra Layer", TECH_GRAY, "MinIO S3 · Meilisearch · BullMQ workers · Docker Compose · GitHub Actions"),
    ]
    for i, (name, col, desc) in enumerate(layers):
        y = Inches(1.55) + i * Inches(0.74)
        add_rect(sl, Inches(0.6), y, Inches(2.2), Inches(0.64), col)
        add_text(sl, name, Inches(0.6), y, Inches(2.2), Inches(0.64),
                 font_size=11, bold=True, color=TECH_BG, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(2.85), y + Inches(0.12), Inches(9.8), Pt(0.8), col)
        add_rect(sl, Inches(2.85), y, Inches(9.8), Inches(0.64), CODE_BG)
        add_text(sl, desc, Inches(3.0), y + Inches(0.1),
                 Inches(9.5), Inches(0.5), font_size=11, color=TECH_WHITE, wrap=True)

    add_text(sl, "Communication inter-layers : HTTP/REST (sync) · BullMQ jobs (async) · SSE streaming (IA) · WebSocket (notif temps réel)",
             Inches(0.5), Inches(6.55), SLIDE_W - Inches(1), Inches(0.4),
             font_size=10.5, color=TECH_ACCENT, italic=True, align=PP_ALIGN.CENTER)

    # ── SLIDE 4 : STACK ──────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "02 — Stack technique détaillée",
                "Choix technologiques et justifications")
    tech_footer(sl, "4 / 14")

    categories = [
        ("BACKEND", TECH_ACCENT, [
            "Node.js 20 LTS + TypeScript 5 (strict: true)",
            "Fastify 4 — performances 3x Express",
            "Drizzle ORM — type-safe, migrations SQL",
            "@fastify/jwt — access + refresh tokens",
            "@fastify/oauth2 — Google, Microsoft",
            "otplib — TOTP MFA",
            "Nodemailer + MJML — emails HTML",
            "PDFKit + @react-pdf/renderer",
            "BullMQ — jobs asynchrones Redis",
        ]),
        ("FRONTEND", TECH_PURPLE, [
            "React 18 + TypeScript 5 + Vite 5",
            "shadcn/ui + Radix UI + Tailwind 3",
            "Framer Motion 11 — animations",
            "Zustand 4 — state management",
            "TanStack Query 5 — data fetching",
            "React Hook Form + Zod — formulaires",
            "React Router 6 — navigation protégée",
            "Recharts — visualisations données",
            "react-i18next — internationalisation",
        ]),
        ("DATA & INFRA", TECH_GREEN, [
            "PostgreSQL 16 — multi-schema tenant",
            "Redis 7 — sessions, cache, queues",
            "MinIO (S3-compatible) — fichiers",
            "Meilisearch — recherche full-text",
            "Docker Compose — dev local",
            "GitHub Actions — CI/CD",
            "AWS SDK v3 — abstraction S3",
        ]),
    ]
    for i, (cat, col, items) in enumerate(categories):
        x = Inches(0.4) + i * Inches(4.25)
        add_rect(sl, x, Inches(1.45), Inches(4.05), Inches(5.0), TECH_PANEL)
        add_rect(sl, x, Inches(1.45), Inches(4.05), Inches(0.07), col)
        add_text(sl, cat, x + Inches(0.15), Inches(1.52),
                 Inches(3.75), Inches(0.5), font_size=13, bold=True, color=col)
        for j, item in enumerate(items):
            add_text(sl, f"  • {item}", x + Inches(0.15), Inches(2.1) + j * Inches(0.5),
                     Inches(3.75), Inches(0.48), font_size=11, color=TECH_WHITE, wrap=True)

    # ── SLIDE 5 : MULTI-TENANT ───────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "03 — Architecture multi-tenant — schema-per-tenant",
                "Stratégie d'isolation : un schéma PostgreSQL dédié par tenant")
    tech_footer(sl, "5 / 14")

    code_block(sl, [
        "-- Schéma global (plateforme)",
        "schema platform : tenants | platform_users | tenant_invitations",
        "",
        "-- Un schéma par tenant (automatiquement provisionné)",
        'schema "tenant_techcorp" : users | employees | departments | contracts',
        '                           payroll_rules | pay_slips | absences | expenses',
        '                           trainings | evaluations | career_skills | ...',
        "",
        "-- Isolation par requête (middleware Fastify)",
        "SET search_path = tenant_techcorp, shared",
        "",
        "-- Jamais de nom de schéma en dur dans le code applicatif",
        "const db = getTenantDbForRequest(request)  // search_path auto-configuré",
    ], Inches(0.5), Inches(1.45), Inches(7.2), Inches(3.5), "PostgreSQL — Isolation multi-schema")

    add_text(sl, "Flux de provisionnement d'un nouveau tenant",
             Inches(8.0), Inches(1.45), Inches(4.9), Inches(0.4),
             font_size=13, bold=True, color=TECH_ACCENT)
    steps = [
        ("1", "INSERT INTO platform.tenants", TECH_ACCENT),
        ("2", "CREATE SCHEMA tenant_{slug}", TECH_GREEN),
        ("3", "Drizzle migrations programmatiques", TECH_PURPLE),
        ("4", "INSERT admin user dans tenant", TECH_YELLOW),
        ("5", "Envoi email invitation (async)", TECH_RED),
    ]
    for i, (num, step, col) in enumerate(steps):
        y = Inches(1.95) + i * Inches(0.6)
        add_rect(sl, Inches(8.0), y, Inches(0.45), Inches(0.5), col)
        add_text(sl, num, Inches(8.0), y, Inches(0.45), Inches(0.5),
                 font_size=14, bold=True, color=TECH_BG, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(8.5), y, Inches(4.4), Inches(0.5), TECH_PANEL)
        add_text(sl, step, Inches(8.65), y + Inches(0.06),
                 Inches(4.1), Inches(0.38), font_size=11, color=TECH_WHITE)

    code_block(sl, [
        "// getTenantDbForRequest — extrait JWT → schemaName",
        "export function getTenantDbForRequest(req: FastifyRequest) {",
        "  const schema = req.user.schemaName",
        "  const client = new Pool({ connectionString: config.db.url })",
        "  client.query(`SET search_path = \"${schema}\", shared`)",
        "  return drizzle(client, { schema: tenantSchema })",
        "}",
    ], Inches(0.5), Inches(5.05), Inches(12.3), Inches(1.5),
    "tenant.ts — résolution automatique par requête")

    # ── SLIDE 6 : AUTH & RBAC ────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "04 — Authentification, autorisation & RBAC",
                "JWT RS256, OAuth2, MFA TOTP, 7 rôles hiérarchiques")
    tech_footer(sl, "6 / 14")

    # Auth flow
    add_rect(sl, Inches(0.5), Inches(1.45), Inches(5.5), Inches(2.4), TECH_PANEL)
    add_text(sl, "🔐  Flux d'authentification", Inches(0.65), Inches(1.5),
             Inches(5.2), Inches(0.4), font_size=13, bold=True, color=TECH_ACCENT)
    auth_steps = [
        "1. POST /auth/login {email, password}",
        "2. findTenantAndUser() → cherche TOUS les tenants candidats",
        "3. bcrypt.compare() → identifie le bon tenant",
        "4. Si MFA activé → challenge TOTP OTP requis",
        "5. Génère accessToken (7j) + refreshToken (30j)",
        "6. Retourne { user, token, tenantConfig } (thème dynamique)",
    ]
    for i, step in enumerate(auth_steps):
        add_text(sl, step, Inches(0.65), Inches(2.0) + i * Inches(0.35),
                 Inches(5.2), Inches(0.33), font_size=10.5, color=TECH_WHITE)

    # Roles
    add_rect(sl, Inches(6.2), Inches(1.45), Inches(6.6), Inches(2.4), TECH_PANEL)
    add_text(sl, "👤  Hiérarchie des rôles", Inches(6.4), Inches(1.5),
             Inches(6.2), Inches(0.4), font_size=13, bold=True, color=TECH_PURPLE)
    roles = [
        ("super_admin", "Plateforme uniquement — zéro accès données RH", TECH_ACCENT),
        ("admin", "Admin tenant — tous droits RH de son tenant", TECH_GREEN),
        ("hr_manager", "RH senior — paie, clôture, recrutement", TECH_GREEN),
        ("hr_officer", "RH opérationnel — saisie & consultation", TECH_YELLOW),
        ("manager", "Manager — son équipe seulement", TECH_YELLOW),
        ("employee", "Self-service — son espace personnel", TECH_GRAY),
        ("readonly", "Lecture seule — aucune modification", TECH_GRAY),
    ]
    for i, (role, desc, col) in enumerate(roles):
        y = Inches(2.0) + i * Inches(0.295)
        add_text(sl, role, Inches(6.4), y, Inches(1.9), Inches(0.28),
                 font_size=10, bold=True, color=col)
        add_text(sl, desc, Inches(8.4), y, Inches(4.2), Inches(0.28),
                 font_size=10, color=TECH_WHITE)

    # Middleware code
    code_block(sl, [
        "// fastify.authenticate — vérifie JWT + set search_path",
        "fastify.decorate('authenticate', async (req, reply) => {",
        "  await req.jwtVerify()",
        "  const { schemaName } = req.user",
        "  await req.server.pg.query(",
        '    `SET search_path = "${schemaName}", shared`',
        "  )",
        "})",
        "",
        "// fastify.authorize(...roles) — contrôle RBAC",
        "fastify.decorate('authorize', (...roles) => async (req, reply) => {",
        "  if (!roles.includes(req.user.role))",
        "    return reply.status(403).send({ error: 'Accès interdit' })",
        "})",
    ], Inches(0.5), Inches(4.0), Inches(5.5), Inches(2.55), "plugins/auth.ts")

    # JWT payload
    code_block(sl, [
        "// Payload JWT",
        "{",
        '  "sub":        "uuid-user-id",',
        '  "tenantId":   "uuid-tenant-id",',
        '  "schemaName": "tenant_techcorp",',
        '  "role":       "hr_manager",',
        '  "email":      "rh@techcorp.com",',
        '  "employeeId": "uuid-emp-id",  // si lié à un employé',
        '  "iat":        1735000000,',
        '  "exp":        1735604800',
        "}",
    ], Inches(6.2), Inches(4.0), Inches(6.6), Inches(2.55), "Structure JWT")

    # ── SLIDE 7 : MODULES API ────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "05 — Modules fonctionnels — API REST",
                "50+ endpoints RESTful, documentés Swagger, protégés RBAC")
    tech_footer(sl, "7 / 14")

    modules_api = [
        ("employees", TECH_ACCENT, [
            "GET    /employees          → liste paginée + filtres",
            "POST   /employees          → créer employé",
            "GET    /employees/:id      → détail complet",
            "PATCH  /employees/:id      → mise à jour",
            "DELETE /employees/:id      → soft delete",
            "POST   /employees/import   → CSV bulk import",
            "GET    /employees/export   → export CSV/XLSX",
        ]),
        ("absences", TECH_GREEN, [
            "GET    /absences           → liste + filtres",
            "POST   /absences           → nouvelle demande",
            "PATCH  /absences/:id/approve → workflow niveau N",
            "PATCH  /absences/:id/reject  → refus motivé",
            "GET    /absences/balances  → soldes par employé",
            "GET    /absences/planning  → planning équipe",
            "GET    /absences/my-absences → self-service",
        ]),
        ("payroll", TECH_PURPLE, [
            "POST   /payroll/calculate  → calcul bulletin",
            "POST   /payroll/close-period → clôture mensuelle",
            "GET    /payroll/payslips   → liste bulletins",
            "GET    /payroll/payslips/:id/pdf → PDF généré",
            "GET    /payroll/my-payslips → self-service",
            "POST   /payroll/dsn        → export DSN URSSAF",
            "GET    /payroll/sepa       → virement SEPA XML",
        ]),
        ("recruitment", TECH_YELLOW, [
            "GET    /recruitment/offers → offres actives",
            "POST   /recruitment/offers → créer une offre",
            "GET    /recruitment/candidates → candidatures",
            "POST   /recruitment/candidates/:id/analyze-cv",
            "POST   /recruitment/offers/:id/source",
            "POST   /recruitment/offers/:id/source/compare",
            "PATCH  /recruitment/candidates/:id/stage",
        ]),
    ]
    for i, (module, col, endpoints) in enumerate(modules_api):
        col_x = Inches(0.4) if i < 2 else Inches(6.8)
        row_y = Inches(1.45) if i % 2 == 0 else Inches(4.15)
        code_block(sl, endpoints, col_x, row_y, Inches(6.0), Inches(2.5), f"/{module}")

    # ── SLIDE 8 : MOTEUR DE PAIE ─────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "06 — Moteur de paie — Logique & calculs",
                "20 rubriques SYNTEC préconfigurées · France 2024")
    tech_footer(sl, "8 / 14")

    code_block(sl, [
        "PayrollEngine.calculate(ctx) : PayrollResult",
        "",
        "ÉTAPE 1 — buildVariables()",
        "  BRUT = salaire brut mensuel",
        "  PLAFOND_SS = 3 864 € (2024)",
        "  TRANCHE_A = min(BRUT, PLAFOND_SS)",
        "  TRANCHE_B = max(0, min(BRUT, 4×PLAFOND_SS) - TRANCHE_A)",
        "  SMIC = 1 766,92 €",
        "  + variableElements (heures supp, primes, IJSS…)",
        "",
        "ÉTAPE 2 — Pour chaque rule (triée par order) :",
        "  safeEval(formula, vars)  →  whitelist [A-Z0-9_\\s\\+\\-\\*\\/\\.\\(\\)]+",
        "  Préfixe VAR:CODE = lire variableElement directement",
        "",
        "ÉTAPE 3 — computeTotals()",
        "  grossSalary  = Σ earning",
        "  employeeCot  = Σ |employee_contribution|",
        "  employerCot  = Σ employer_contribution",
        "  netBeforeTax = grossSalary - employeeCot",
        "  netPayable   = netBeforeTax  // PAS via DSN ultérieurement",
        "  employerCost = grossSalary + employerCot",
    ], Inches(0.5), Inches(1.45), Inches(7.0), Inches(5.1), "PayrollEngine — Algorithme")

    add_text(sl, "Rubriques SYNTEC préconfigurées (20)",
             Inches(7.7), Inches(1.45), Inches(5.1), Inches(0.4),
             font_size=13, bold=True, color=TECH_ACCENT)
    rubriques = [
        ("1000", "Salaire de base", "earning"),
        ("2000", "Congés payés", "earning"),
        ("4100", "CSG déductible", "emp 6,80%"),
        ("4110", "CSG non déductible", "emp 2,40%"),
        ("4120", "CRDS", "emp 0,50%"),
        ("4210", "Maladie patronale", "pat 7,00%"),
        ("4300", "Alloc. familiales", "pat 3,45%"),
        ("4400", "AT/MP", "pat 2,22%"),
        ("4500", "Retraite base sal.", "emp 6,90%"),
        ("4510", "Retraite base pat.", "pat 8,55%"),
        ("4600", "AGIRC-ARRCO T1 sal.", "emp 3,15%"),
        ("4610", "AGIRC-ARRCO T1 pat.", "pat 4,72%"),
        ("5000", "Mutuelle sal.", "45 €"),
        ("5010", "Mutuelle pat.", "90 €"),
    ]
    for i, (code, label, rate) in enumerate(rubriques):
        y = Inches(1.92) + i * Inches(0.36)
        col_bg = CODE_BG if i % 2 == 0 else TECH_PANEL
        add_rect(sl, Inches(7.7), y, Inches(5.1), Inches(0.35), col_bg)
        add_text(sl, code, Inches(7.75), y + Inches(0.03),
                 Inches(0.6), Inches(0.3), font_size=9.5, color=TECH_ACCENT,
                 bold=True)
        add_text(sl, label, Inches(8.4), y + Inches(0.03),
                 Inches(2.8), Inches(0.3), font_size=9.5, color=TECH_WHITE)
        add_text(sl, rate, Inches(11.25), y + Inches(0.03),
                 Inches(1.4), Inches(0.3), font_size=9.5, color=TECH_YELLOW,
                 align=PP_ALIGN.RIGHT)

    # ── SLIDE 9 : IA INTÉGRATION ─────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "07 — Intégration IA — Claude & Mistral",
                "SDK natif · SSE streaming · Comparateur automatique · Scoring qualité")
    tech_footer(sl, "9 / 14")

    code_block(sl, [
        "// CLAUDE — SDK @anthropic-ai/sdk",
        "const client = new Anthropic({ apiKey: process.env.ANTHROPIC_API_KEY })",
        "",
        "async function callClaude(prompt: string): Promise<AIResponse> {",
        "  const t0 = Date.now()",
        "  const msg = await client.messages.create({",
        "    model:      process.env.AI_MODEL ?? 'claude-sonnet-4-20250514',",
        "    max_tokens: 2000,",
        "    messages:   [{ role: 'user', content: prompt }]",
        "  })",
        "  const text = msg.content[0]?.type === 'text' ? msg.content[0].text : ''",
        "  return {",
        "    provider: 'claude', text,",
        "    inputTokens:  msg.usage.input_tokens,",
        "    outputTokens: msg.usage.output_tokens,",
        "    latencyMs:    Date.now() - t0,",
        "    estimatedCostEur: (msg.usage.input_tokens * 3 + msg.usage.output_tokens * 15) / 1e6 * 0.92",
        "  }",
        "}",
    ], Inches(0.5), Inches(1.45), Inches(6.2), Inches(3.8), "Claude API Integration")

    code_block(sl, [
        "// MISTRAL — API REST native (fetch)",
        "async function callMistral(prompt: string): Promise<AIResponse> {",
        "  const t0 = Date.now()",
        "  const res = await fetch('https://api.mistral.ai/v1/chat/completions', {",
        "    method: 'POST',",
        "    headers: { Authorization: `Bearer ${process.env.MISTRAL_API_KEY}` },",
        "    body: JSON.stringify({",
        "      model:    process.env.MISTRAL_MODEL ?? 'mistral-large-latest',",
        "      messages: [{ role: 'user', content: prompt }]",
        "    })",
        "  })",
        "  const data = await res.json()",
        "  return {",
        "    provider: 'mistral',",
        "    text:         data.choices[0].message.content,",
        "    inputTokens:  data.usage.prompt_tokens,",
        "    outputTokens: data.usage.completion_tokens,",
        "    latencyMs:    Date.now() - t0,",
        "    estimatedCostEur: ...",
        "  }",
        "}",
    ], Inches(6.9), Inches(1.45), Inches(6.2), Inches(3.8), "Mistral API Integration")

    code_block(sl, [
        "// Comparateur automatique — POST /recruitment/offers/:id/source/compare",
        "const [claudeResult, mistralResult] = await Promise.allSettled([",
        "  callClaude(buildSourcingPrompt(offer, platforms, maxProfiles)),",
        "  callMistral(buildSourcingPrompt(offer, platforms, maxProfiles))",
        "])",
        "// computeRichnessScore() → score objectif basé sur : nb profils, champs remplis, qualité JSON",
        "// Recommendation automatique : winner = meilleur score | latence | coût",
    ], Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.15), "Endpoint comparateur — /source/compare")

    # ── SLIDE 10 : SÉCURITÉ & RGPD ───────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "10 — Sécurité & conformité RGPD",
                "Chiffrement, isolation, audit trail, durées de conservation")
    tech_footer(sl, "10 / 14")

    security_areas = [
        ("🔒  Chiffrement des données sensibles", TECH_RED, [
            "NIR (numéro sécu) → AES-256-GCM chiffré en base",
            "IBAN → AES-256-GCM chiffré en base",
            "Mots de passe → bcrypt (cost 12)",
            "JWT → HS256 (secret 64+ chars) ou RS256",
            "Fichiers S3 → chiffrement côté serveur (SSE-S3)",
        ]),
        ("🛡  Isolation multi-tenant", TECH_ACCENT, [
            "Schéma PostgreSQL dédié par tenant",
            "search_path forcé sur CHAQUE requête DB",
            "Zéro requête cross-schema possible en applicatif",
            "super_admin : accès platform uniquement",
            "Audit log cross-tenant impossible par design",
        ]),
        ("📋  Conformité RGPD", TECH_GREEN, [
            "Droit à l'oubli → DELETE + purge S3 automatisée",
            "Export RGPD → /platform/tenants/:id/export",
            "Durées légales : bulletins 50 ans | dossier 5 ans",
            "Logs audit immuables (insert only, pas d'update)",
            "DPA disponible sur demande",
        ]),
        ("🔐  Contrôles d'accès", TECH_PURPLE, [
            "RBAC 7 rôles + middleware authorize() Fastify",
            "Rate limiting : 100 req/min par IP (fastify/rate-limit)",
            "CORS strict : origines whitelistées uniquement",
            "MFA TOTP obligatoire (configurable par tenant)",
            "Tokens refresh rotation à chaque usage",
        ]),
    ]
    for i, (title, col, items) in enumerate(security_areas):
        col_x = Inches(0.4) if i < 2 else Inches(6.8)
        row_y = Inches(1.45) if i % 2 == 0 else Inches(4.15)
        add_rect(sl, col_x, row_y, Inches(6.0), Inches(2.45), TECH_PANEL)
        add_rect(sl, col_x, row_y, Inches(6.0), Inches(0.07), col)
        add_text(sl, title, col_x + Inches(0.15), row_y + Inches(0.12),
                 Inches(5.7), Inches(0.5), font_size=13, bold=True, color=col)
        for j, item in enumerate(items):
            add_text(sl, f"  • {item}", col_x + Inches(0.15), row_y + Inches(0.7) + j * Inches(0.36),
                     Inches(5.7), Inches(0.34), font_size=11, color=TECH_WHITE, wrap=True)

    # ── SLIDE 11 : INFRASTRUCTURE ────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "11 — Infrastructure & déploiement",
                "Docker Compose (dev) · Docker multi-stage (prod) · GitHub Actions CI/CD")
    tech_footer(sl, "11 / 14")

    code_block(sl, [
        "# docker-compose.yml — Services dev",
        "services:",
        "  postgres:",
        "    image: postgres:16-alpine",
        "    ports: ['5433:5432']",
        "    environment: { POSTGRES_DB: nexusrh, POSTGRES_USER: nexusrh }",
        "",
        "  redis:",
        "    image: redis:7-alpine",
        "    ports: ['6379:6379']",
        "",
        "  meilisearch:",
        "    image: getmeili/meilisearch:v1.6",
        "    ports: ['7700:7700']",
        "",
        "  minio:",
        "    image: minio/minio",
        "    ports: ['9000:9000', '9001:9001']",
        "    command: server /data --console-address ':9001'",
    ], Inches(0.5), Inches(1.45), Inches(5.8), Inches(5.1), "Infrastructure locale (dev)")

    add_text(sl, "Déploiement production", Inches(6.6), Inches(1.45),
             Inches(6.2), Inches(0.4), font_size=13, bold=True, color=TECH_ACCENT)
    prod_items = [
        ("🐳  Docker multi-stage", "Images optimisées : api (150MB) + web (25MB Nginx static)"),
        ("⚙  Variables d'env", "Secrets injectés via Docker secrets ou Vault"),
        ("🔄  CI/CD", "GitHub Actions : lint → test → build → push registry → deploy"),
        ("📊  Monitoring", "Pino logs (JSON) → Grafana Loki · métriques → Prometheus"),
        ("🔁  Backup", "pg_dump automatique toutes les 6h → S3 chiffré"),
        ("📈  Scalabilité", "API stateless → scale horizontal · BullMQ workers séparés"),
        ("🌐  Reverse proxy", "Nginx / Traefik · TLS Let's Encrypt · HTTP/2"),
        ("🔒  Secrets", "JWT_SECRET ≥ 64 chars · rotation automatique refresh tokens"),
    ]
    for i, (title, desc) in enumerate(prod_items):
        y = Inches(1.95) + i * Inches(0.65)
        add_rect(sl, Inches(6.6), y, Inches(6.2), Inches(0.58), TECH_PANEL)
        add_rect(sl, Inches(6.6), y, Inches(0.06), Inches(0.58), TECH_ACCENT)
        add_text(sl, title, Inches(6.8), y + Inches(0.06),
                 Inches(2.3), Inches(0.45), font_size=11, bold=True, color=TECH_ACCENT)
        add_text(sl, desc, Inches(9.2), y + Inches(0.06),
                 Inches(3.5), Inches(0.45), font_size=10.5, color=TECH_WHITE, wrap=True)

    # ── SLIDE 12 : PERFORMANCE ───────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "12 — Performance & scalabilité",
                "Benchmarks, optimisations, SLAs cibles")
    tech_footer(sl, "12 / 14")

    kpis_perf = [
        ("< 50ms", "P50 latence API\n(hors calcul paie)", TECH_ACCENT),
        ("< 200ms", "P99 latence API\n(avec cache Redis)", TECH_GREEN),
        ("10 000+", "Requêtes/sec\n(Fastify, single node)", TECH_PURPLE),
        ("99,9%", "SLAs uptime cible\n(prod enterprise)", TECH_YELLOW),
    ]
    for i, (num, label, col) in enumerate(kpis_perf):
        x = Inches(0.4) + i * Inches(3.2)
        add_rect(sl, x, Inches(1.45), Inches(3.0), Inches(1.8), TECH_PANEL)
        add_rect(sl, x, Inches(1.45), Inches(3.0), Inches(0.07), col)
        add_text(sl, num, x + Inches(0.15), Inches(1.6),
                 Inches(2.7), Inches(0.75), font_size=30, bold=True, color=col)
        add_text(sl, label, x + Inches(0.15), Inches(2.35),
                 Inches(2.7), Inches(0.75), font_size=11, color=TECH_WHITE, wrap=True)

    add_text(sl, "Stratégies d'optimisation", Inches(0.5), Inches(3.5),
             Inches(12), Inches(0.45), font_size=15, bold=True, color=TECH_WHITE)
    opts = [
        ("🗃  Connection pooling", "PgBouncer + pool Drizzle (min 2, max 20) — connexions DB mutualisées"),
        ("⚡  Cache Redis", "Sessions JWT · résultats calcul KPIs · soldes congés (TTL 5 min)"),
        ("🔍  Index DB", "Tous les FK indexés · index composites email+tenantId · GIN sur JSONB"),
        ("📄  Lazy migrations", "ensureTenantSchema() avec Set<string> en mémoire — exécution unique par schéma"),
        ("📦  Code splitting", "Vite chunks par route · lazy imports · prefetch sur hover"),
        ("🌊  Streaming SSE", "Réponses IA streamées → UX instantanée sans bloquer l'event loop"),
        ("🔄  BullMQ workers", "Jobs lourds (paie batch, scoring IA, emails) hors du thread HTTP"),
    ]
    for i, (title, desc) in enumerate(opts):
        y = Inches(4.05) + i * Inches(0.45)
        add_rect(sl, Inches(0.5), y, SLIDE_W - Inches(1), Inches(0.4), TECH_PANEL)
        add_rect(sl, Inches(0.5), y, Inches(0.06), Inches(0.4), TECH_ACCENT)
        add_text(sl, title, Inches(0.7), y + Inches(0.04),
                 Inches(2.8), Inches(0.32), font_size=11, bold=True, color=TECH_ACCENT)
        add_text(sl, desc, Inches(3.6), y + Inches(0.04),
                 Inches(9.4), Inches(0.32), font_size=11, color=TECH_WHITE, wrap=True)

    # ── SLIDE 13 : TESTS ─────────────────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    tech_header(sl, "13 — Tests & qualité",
                "Vitest (unit) · Playwright (E2E) · TypeScript strict · CI obligatoire")
    tech_footer(sl, "13 / 14")

    code_block(sl, [
        "// Tests unitaires — Vitest",
        "describe('PayrollEngine', () => {",
        "  it('calcule net >= 0 pour SMIC', () => {",
        "    const result = PayrollEngine.calculate({",
        "      grossSalary: 1766.92, rules: baseRules, ...",
        "    })",
        "    expect(result.netPayable).toBeGreaterThan(0)",
        "    expect(result.netPayable).toBeLessThan(result.grossSalary)",
        "  })",
        "})",
        "",
        "// Tests E2E — Playwright",
        "test('employee voit ses bulletins', async ({ page }) => {",
        "  await page.goto('/login')",
        "  await page.fill('[name=email]', 'employe@techcorp.com')",
        "  await page.fill('[name=password]', 'Admin1234!')",
        "  await page.click('[type=submit]')",
        "  await page.goto('/mon-espace/bulletins')",
        "  await expect(page.locator('.payslip-card')).toHaveCount(6)",
        "})",
    ], Inches(0.5), Inches(1.45), Inches(7.0), Inches(4.8), "Tests automatisés")

    add_text(sl, "Couverture & CI", Inches(7.7), Inches(1.45),
             Inches(5.1), Inches(0.4), font_size=13, bold=True, color=TECH_ACCENT)
    ci_items = [
        ("Vitest unit", "Moteur paie, auth, RBAC, provisionnement"),
        ("Playwright E2E", "Auth, self-service employé, création tenant, workflow absences"),
        ("TypeScript strict", "noUncheckedIndexedAccess, strictNullChecks — zéro any"),
        ("ESLint + Prettier", "Qualité code homogène sur tout le monorepo"),
        ("CI GitHub Actions", "lint → tsc → test → build à chaque PR"),
        ("Coverage cible", "70% unit · 90% E2E parcours critiques"),
        ("Security audit", "npm audit + Dependabot alerts activés"),
    ]
    for i, (tool, desc) in enumerate(ci_items):
        y = Inches(1.95) + i * Inches(0.62)
        add_rect(sl, Inches(7.7), y, Inches(5.1), Inches(0.55), TECH_PANEL)
        add_rect(sl, Inches(7.7), y, Inches(0.06), Inches(0.55), TECH_GREEN)
        add_text(sl, tool, Inches(7.9), y + Inches(0.07),
                 Inches(1.8), Inches(0.4), font_size=11, bold=True, color=TECH_GREEN)
        add_text(sl, desc, Inches(9.8), y + Inches(0.07),
                 Inches(2.9), Inches(0.4), font_size=10.5, color=TECH_WHITE, wrap=True)

    # ── SLIDE 14 : ANNEXES / CLÔTURE ─────────────────────────────────────────
    sl = blank_slide(prs)
    fill_bg(sl, TECH_BG)
    add_rect(sl, 0, 0, Inches(0.06), SLIDE_H, TECH_ACCENT)
    add_rect(sl, Inches(0.5), SLIDE_H / 2 - Pt(1), SLIDE_W - Inches(0.5), Pt(2), TECH_ACCENT)

    add_text(sl, "Annexes & Ressources",
             Inches(1.0), Inches(0.8), Inches(11), Inches(0.9),
             font_size=36, bold=True, color=TECH_WHITE, align=PP_ALIGN.CENTER)

    annexes = [
        ("📄  Swagger API", "http://localhost:4000/docs — Documentation interactive complète de tous les endpoints"),
        ("🗄  Schéma DB", "apps/api/src/db/schema/ — Drizzle schema TypeScript (platform + tenant)"),
        ("🌱  Seed données", "pnpm --filter api run db:seed — 2 tenants, 68 employés, 300 bulletins"),
        ("🔧  Variables env", ".env.example — toutes les variables documentées avec valeurs par défaut"),
        ("🐳  Docker", "docker-compose.yml — 5 services : postgres, redis, meilisearch, minio, mailhog"),
        ("📝  CLAUDE.md", "/CLAUDE.md — Document maître : architecture, patterns, pièges en production"),
        ("🧪  Tests", "pnpm test — unit (Vitest) · pnpm e2e — parcours complets (Playwright)"),
        ("📊  Métriques", "GET /health — healthcheck · GET /metrics — Prometheus (si activé)"),
    ]
    for i, (title, desc) in enumerate(annexes):
        col = 0 if i < 4 else 1
        row = i if i < 4 else i - 4
        x = Inches(0.8) + col * Inches(6.4)
        y = Inches(1.9) + row * Inches(0.88)
        add_rect(sl, x, y, Inches(6.0), Inches(0.75), TECH_PANEL)
        add_rect(sl, x, y, Inches(0.06), Inches(0.75), TECH_ACCENT)
        add_text(sl, title, x + Inches(0.18), y + Inches(0.06),
                 Inches(5.7), Inches(0.3), font_size=11, bold=True, color=TECH_ACCENT)
        add_text(sl, desc, x + Inches(0.18), y + Inches(0.38),
                 Inches(5.7), Inches(0.3), font_size=10, color=TECH_WHITE, wrap=True)

    add_text(sl, "© 2025 OpenLab — NexusRH Spécifications Fonctionnelles v1.0  |  Toute reproduction interdite sans accord.",
             Inches(0.8), Inches(5.85), SLIDE_W - Inches(1.6), Inches(0.35),
             font_size=10, color=TECH_GRAY, italic=True, align=PP_ALIGN.CENTER)
    tech_footer(sl, "14 / 14")

    prs.save(output_path)
    print(f"[OK] Document technique sauvegarde : {output_path}")


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    out_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "docs")
    os.makedirs(out_dir, exist_ok=True)

    commercial_path = os.path.join(out_dir, "NexusRH_Presentation_Commerciale.pptx")
    technical_path  = os.path.join(out_dir, "NexusRH_Specifications_Fonctionnelles.pptx")

    print("\nGeneration des documents NexusRH...")
    build_commercial_deck(commercial_path)
    build_technical_spec(technical_path)
    print(f"\nDossier : {out_dir}")
    print("Termine.\n")
